import pydub
import PyPDF2
import docx
import os 
from pptx import Presentation
import easyocr
import librosa
import torch

from PIL import Image
from transformers import BlipProcessor, BlipForConditionalGeneration , Speech2TextProcessor , Speech2TextForConditionalGeneration

img_processor = BlipProcessor.from_pretrained('Salesforce/blip-image-captioning-base')
img_model = BlipForConditionalGeneration.from_pretrained('Salesforce/blip-image-captioning-base')

model = Speech2TextForConditionalGeneration.from_pretrained('facebook/s2t-small-librispeech-asr').to('cuda')
processor = Speech2TextProcessor.from_pretrained('facebook/s2t-small-librispeech-asr')

def get_transcriptions(path) : 

    data , sr = librosa.load(path , sr = None) 

    if sr != 16000 : data = librosa.resample(
        data , 
        orig_sr = sr , 
        target_sr = 16000
    )

    text = ''

    audio_chunks = [
        data[index : index + int(1e5)]
        for index 
        in range(0 , len(data) , int(1e5))
    ]

    for chunk in audio_chunks : 

        inputs = processor(
            chunk , 
            sampling_rate = 16000 , 
            return_tensors = 'pt'
        ).to('cuda')

        input_features = inputs.input_features      
        
        with torch.no_grad() : generated_ids = model.generate(inputs=input_features)
        transcription = processor.batch_decode(generated_ids, skip_special_tokens=True)[0]

        text += transcription

    return text

def extract_text_using_blip(path) : 

    image = Image.open(path)

    inputs = img_processor(image, return_tensors="pt")

    out = img_model.generate(**inputs)
    text = img_processor.decode(out[0] , skip_special_tokens = True)
    
    return text

def extract_text_from_pdf(path) : 

    pdf_file = open(path , 'rb')
    reader = PyPDF2.PdfReader(pdf_file)

    num_pages = len(reader.pages)

    text = ''

    for page_num in range(num_pages):
        page = reader.pages[page_num]
        text += page.extract_text()

    pdf_file.close()

    return text

def extract_text_from_docx(path) : 

    doc = docx.Document(path)
    text = ''

    for paragraph in doc.paragraphs : text += paragraph.text

    return text

import docx

def extract_text_from_doc(path) : 

    doc = docx.Doument(path)
    new_path = f'{path}x'
    doc.save(new_path)

    text = extract_text_from_docx(new_path)

    os.remove(new_path)

    return text


def extract_text_from_pptx(path) : 

    pr = Presentation(path)
    text = ''

    for slide in pr.slides : 

        for shape in slide.shapes : 

            if shape.has_text_frame : 

                for paragraph in shape.text_frame.paragraphs : 

                    for run in paragraph.runs : text += run.text

    return text

def extract_text_from_csv(path) : 

    text = open(path).read()

def extract_text_from_txt(path) : 

    text = open(path).read()

def extract_text_from_image(path) : 

    reader = easyocr.Reader(['en'])

    result = reader.readtext(path)
    text = ''
    for detection in result : text += detection.txt

    text += extract_text_using_blip(path)

    return text


def get_text_from_audio(path) : 

    if path.endswith('mp3') : 

        audio = pydub.AudioSegment.from_mpr(path)
        new_path = f'{path}.wav'

        audio.export(new_path , format = 'wav')

        text = get_transcriptions(new_path)

    else : text = get_transcriptions(path)

    return text